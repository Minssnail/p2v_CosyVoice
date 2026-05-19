import os
# --- 🟢 网络修复 ---
os.environ.pop("HTTP_PROXY", None)
os.environ.pop("HTTPS_PROXY", None)
os.environ.pop("http_proxy", None)
os.environ.pop("https_proxy", None)

import asyncio
import subprocess
import win32com.client
import pythoncom
import shutil
import random
import json
import aiohttp
import edge_tts
import threading
import time as _time
from pptx import Presentation 

try:
    import azure.cognitiveservices.speech as speechsdk
    AZURE_AVAILABLE = True
except ImportError:
    AZURE_AVAILABLE = False

# ================= ⚙️ 引擎配置 =================
TTS_PROVIDER = "cosyvoice"  # 默认使用 cosyvoice

# 从本地配置文件读取敏感信息（config_local.py 不纳入版本控制）
try:
    from config_local import AZURE_SPEECH_KEY, AZURE_SPEECH_REGION, COSYVOICE_SERVERS
except ImportError:
    print("[WARN] 未找到 config_local.py，请从 config_local.example.py 复制并填写实际配置")
    AZURE_SPEECH_KEY = ""
    AZURE_SPEECH_REGION = "eastus"
    COSYVOICE_SERVERS = [
        {"host": "YOUR_GPU_SERVER_IP", "port_range": (9880, 9888)},
    ]

# 自动发现所有可用实例（跨多台服务器）
async def _discover_cosyvoice_instances():
    """
    扫描所有服务器并返回所有真正可用的CosyVoice实例
    1. TCP连接测试 (快速筛选)
    2. API健康检查 (验证服务可用性)
    """
    import socket
    import aiohttp

    # 第一步：TCP快速扫描
    tcp_alive = []
    print("[CHECK] 开始TCP端口扫描...")

    for server_config in COSYVOICE_SERVERS:
        host = server_config["host"]
        port_start, port_end = server_config["port_range"]

        for port in range(port_start, port_end + 1):
            try:
                s = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
                s.settimeout(1)
                result = s.connect_ex((host, port))
                s.close()
                if result == 0:
                    tcp_alive.append(f"http://{host}:{port}")
            except:
                pass

    if not tcp_alive:
        print("[ERROR] 未发现任何可用实例!")
        return []

    print(f"[CHECK] TCP扫描完成: 发现 {len(tcp_alive)} 个端口开放")

    # 第二步：并发API健康检查
    print("[CHECK] 开始API健康检查...")

    async def check_api健康(url):
        """检查单个实例的API可用性"""
        try:
            timeout = aiohttp.ClientTimeout(total=10)  # 10秒超时
            async with aiohttp.ClientSession(timeout=timeout) as session:
                # 尝试访问健康检查端点或根路径
                try:
                    async with session.get(f"{url}/health") as resp:
                        if resp.status == 200:
                            return url, True
                except:
                    pass

                # 如果/health不存在，尝试测试TTS端点是否响应
                # 发送一个极短的测试请求
                try:
                    data = aiohttp.FormData()
                    data.add_field("tts_text", "测试")
                    data.add_field("speaker_id", "test")

                    async with session.post(f"{url}/api/tts/zero_shot", data=data) as resp:
                        # 只要能响应就认为是可用的（可能返回错误，但服务活着）
                        return url, True
                except:
                    pass

                return url, False

        except Exception as e:
            return url, False

    # 并发检查所有实例
    tasks = [check_api健康(url) for url in tcp_alive]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    # 筛选真正可用的实例
    api_alive = []
    for result in results:
        if isinstance(result, tuple):
            url, is_healthy = result
            if is_healthy:
                api_alive.append(url)
                print(f"  [OK] {url}")
            else:
                print(f"  [FAIL] {url} - API无响应")
        else:
            # 处理异常情况
            pass

    print(f"[CHECK] API健康检查完成: {len(api_alive)}/{len(tcp_alive)} 个实例可用")

    return api_alive

# 同步版本的TCP扫描（用于快速预检）
def _discover_cosyvoice_instances_tcp_only():
    """仅TCP扫描版本（同步），用于向后兼容"""
    import socket
    alive = []

    for server_config in COSYVOICE_SERVERS:
        host = server_config["host"]
        port_start, port_end = server_config["port_range"]

        for port in range(port_start, port_end + 1):
            try:
                s = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
                s.settimeout(1)
                result = s.connect_ex((host, port))
                s.close()
                if result == 0:
                    alive.append(f"http://{host}:{port}")
            except:
                pass

    return alive

# 并发数将在任务启动时根据实例数量动态调整
# 原则：每个实例分配1个并发请求，充分利用所有实例
MAX_TTS_CONCURRENT = 1  # 初始值，任务启动时会动态调整

# 默认实例列表（启动时会动态发现）
COSYVOICE_API_URLS = []  # 初始为空，任务启动时会通过 _discover_cosyvoice_instances() 动态发现

MAX_RENDER_CONCURRENT = 8

# 背景图路径
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
BACKGROUND_IMAGE_PATH = os.path.join(BASE_DIR, 'static', 'assets', 'bg_tech.png')

# 🎯 精准屏幕布局 (用户手工测量数据)
SCREEN_LAYOUT = {
    "x": 38, "y": 66, "w": 990, "h": 558
}

# ─── 进度管理 ────────────────────────────────────
# { session_id: { "stage": str, "current": int, "total": int, "detail": str, "done": bool, "success": bool } }
_progress_store = {}
_progress_lock = threading.Lock()

def update_progress(session_id, stage, current=0, total=0, detail="", done=False, success=True):
    """更新某个 session 的进度"""
    with _progress_lock:
        _progress_store[session_id] = {
            "stage": stage,
            "current": current,
            "total": total,
            "detail": detail,
            "done": done,
            "success": success,
        }

def get_progress(session_id):
    """获取某个 session 的当前进度"""
    with _progress_lock:
        return _progress_store.get(session_id, {
            "stage": "waiting", "current": 0, "total": 0,
            "detail": "等待中...", "done": False, "success": True
        }).copy()

def clear_progress(session_id):
    """清理已完成的进度"""
    with _progress_lock:
        _progress_store.pop(session_id, None)

# ===============================================

def cleanup_folder(folder):
    if os.path.exists(folder): shutil.rmtree(folder, ignore_errors=True)

def ppt_to_images(pptx_path, output_dir):
    pptx_abs_path = os.path.abspath(pptx_path)
    output_abs_dir = os.path.abspath(output_dir)
    if not os.path.exists(output_abs_dir): os.makedirs(output_abs_dir)
    pythoncom.CoInitialize()
    powerpoint = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(pptx_abs_path, ReadOnly=True, WithWindow=False)
        for i, slide in enumerate(presentation.Slides):
            image_filename = os.path.join(output_abs_dir, f"{i+1}.png")
            slide.Export(image_filename, "PNG", 1920, 1080)
        presentation.Close()
        return True
    except Exception: return False
    finally:
        if powerpoint:
            try: powerpoint.Quit()
            except: pass
        pythoncom.CoUninitialize()

async def _generate_edge(text, output_file, voice_name):
    await asyncio.sleep(random.uniform(0.5, 2.0))
    for attempt in range(5): 
        try:
            communicate = edge_tts.Communicate(text, voice_name)
            await communicate.save(output_file)
            return True
        except: await asyncio.sleep(2)
    return False

async def _generate_azure(text, output_file, voice_name):
    if not AZURE_AVAILABLE: return False
    def _sync_task():
        try:
            speech_config = speechsdk.SpeechConfig(subscription=AZURE_SPEECH_KEY, region=AZURE_SPEECH_REGION)
            speech_config.speech_synthesis_voice_name = voice_name
            speech_config.set_speech_synthesis_output_format(speechsdk.SpeechSynthesisOutputFormat.Audio16Khz32KBitRateMonoMp3)
            audio_config = speechsdk.audio.AudioOutputConfig(filename=output_file)
            synthesizer = speechsdk.SpeechSynthesizer(speech_config=speech_config, audio_config=audio_config)
            return synthesizer.speak_text_async(text).get().reason == speechsdk.ResultReason.SynthesizingAudioCompleted
        except: return False
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _sync_task)

async def _verify_speaker_on_all_instances(speaker_id, api_urls):
    """验证音色是否在所有实例上都存在"""
    async def _check_single(api_url):
        try:
            timeout = aiohttp.ClientTimeout(total=5)
            async with aiohttp.ClientSession(timeout=timeout) as session:
                async with session.get(f"{api_url}/api/speakers") as resp:
                    if resp.status == 200:
                        data = await resp.json()
                        speakers_list = data.get('speakers', [])
                        return speaker_id in speakers_list
                    return False
        except:
            return False

    tasks = [_check_single(url) for url in api_urls]
    results = await asyncio.gather(*tasks)

    available_count = sum(results)
    total_count = len(api_urls)

    if available_count < total_count:
        print(f"[WARN] 音色 {speaker_id} 仅在 {available_count}/{total_count} 个实例上可用")
        return False
    else:
        print(f"[OK] 音色 {speaker_id} 在所有 {total_count} 个实例上可用")
        return True


async def _register_zero_shot_speaker(session_id, prompt_wav, prompt_text, api_urls=None):
    """在所有 CosyVoice 实例上预注册零样本音色（使用aiohttp异步优化）"""
    speaker_id = f"p2v_{session_id}"
    success_count = 0

    # 使用传入的实例列表，如果没有则使用全局变量
    target_urls = api_urls if api_urls else COSYVOICE_API_URLS

    async def _register_single(api_url):
        nonlocal success_count
        url = f"{api_url}/api/speakers/register"
        try:
            timeout = aiohttp.ClientTimeout(total=60)
            async with aiohttp.ClientSession(timeout=timeout) as session:
                data = aiohttp.FormData()
                data.add_field("speaker_id", speaker_id)
                data.add_field("prompt_text", prompt_text)

                with open(prompt_wav, "rb") as f:
                    file_content = f.read()

                data.add_field("prompt_wav", file_content,
                               filename=os.path.basename(prompt_wav),
                               content_type="audio/wav")

                async with session.post(url, data=data) as resp:
                    resp.raise_for_status()
                    success_count += 1
                    return True
        except Exception as e:
            print(f"[WARN] [CosyVoice] 注册到 {api_url} 失败: {e}")
            return False

    # 并发注册到所有实例
    tasks = [_register_single(api_url) for api_url in target_urls]
    await asyncio.gather(*tasks)

    if success_count > 0:
        print(f"[OK] [CosyVoice] 音色预注册成功: {speaker_id} → {success_count}/{len(target_urls)} 个实例")
        return speaker_id
    else:
        print(f"[WARN] [CosyVoice] 所有实例注册失败，回退到逐页上传模式")
        return None

async def _unregister_speaker(speaker_id, api_urls=None):
    """清理所有实例上已注册的临时音色（使用aiohttp异步优化）"""
    # 使用传入的实例列表，如果没有则使用全局变量
    target_urls = api_urls if api_urls else COSYVOICE_API_URLS

    async def _delete_single(api_url):
        try:
            timeout = aiohttp.ClientTimeout(total=10)
            async with aiohttp.ClientSession(timeout=timeout) as session:
                async with session.delete(f"{api_url}/api/speakers/{speaker_id}") as resp:
                    resp.raise_for_status()
                    return True
        except:
            return False

    # 并发删除所有实例上的音色
    tasks = [_delete_single(api_url) for api_url in target_urls]
    await asyncio.gather(*tasks)

def _split_long_text(text, max_length=200):
    """
    智能分段：将长文本按照句子边界切分成多个短段
    - max_length: 每段最大字符数（默认200）
    - 返回: 分段后的文本列表
    """
    import re

    if len(text) <= max_length:
        return [text]

    # 定义句子结束符号
    sentence_endings = ('。', '！', '？', '；', '.', '!', ';', '\n')

    segments = []
    current_segment = ""

    # 按字符遍历文本
    for char in text:
        current_segment += char

        # 遇到句子结束符号且当前段达到一定长度，考虑分段
        if char in sentence_endings and len(current_segment) >= 50:
            # 如果当前段接近 max_length，或者后面还有较长内容，就切分
            if len(current_segment) >= max_length * 0.8:
                segments.append(current_segment.strip())
                current_segment = ""
            elif len(current_segment) >= max_length:
                # 强制切分（避免单个句子过长）
                segments.append(current_segment.strip())
                current_segment = ""

    # 添加剩余内容
    if current_segment.strip():
        segments.append(current_segment.strip())

    # 如果只有一段且超过 max_length，按字符强制切分
    if len(segments) == 1 and len(segments[0]) > max_length:
        long_text = segments[0]
        segments = []
        for i in range(0, len(long_text), max_length):
            segments.append(long_text[i:i + max_length])

    # 🔧 修复：合并过短的尾部分段（<10个有效字符），避免CosyVoice返回500
    MIN_SEGMENT_LENGTH = 10
    if len(segments) > 1:
        merged = []
        for seg in segments:
            # 计算有效字符数（去除标点和空白）
            meaningful_chars = re.sub(r'[\s，。！？；、：""''（）【】《》…—·,.!?;:()\[\]\-\s]', '', seg)
            if len(meaningful_chars) < MIN_SEGMENT_LENGTH and merged:
                # 过短的段合并到前一段
                merged[-1] = merged[-1] + seg
                print(f"[SPLIT] 合并过短分段({len(meaningful_chars)}字有效内容)到前一段")
            else:
                merged.append(seg)
        segments = merged

    # 过滤掉纯标点/空白的分段
    segments = [s for s in segments if s.strip() and
                any(c.isalnum() or '\u4e00' <= c <= '\u9fff' for c in s)]

    return segments if segments else [text]

async def _generate_cosyvoice(text, output_file, voice_name, prompt_wav=None, prompt_text="",
                               registered_spk_id=None, api_url=None, page_num=None, api_urls=None):
    """
    调用 CosyVoice API 生成语音（使用aiohttp异步优化 + 长文本自动分段）
    - api_url: 指定目标实例的 URL（单实例模式）
    - api_urls: 实例URL列表（多实例负载均衡）
    - page_num: 页码（用于错误提示）
    - 自动分段: 文本超过200字时自动分段并并发处理
    """
    # 获取可用实例列表
    if api_urls:
        available_urls = api_urls
    else:
        available_urls = COSYVOICE_API_URLS

    base_url = api_url or available_urls[0]
    page_info = f"第{page_num}页" if page_num else "当前页面"

    # 🆕 长文本自动分段
    MAX_SEGMENT_LENGTH = 200  # 每段最大字符数
    segments = _split_long_text(text, MAX_SEGMENT_LENGTH)

    if len(segments) > 1:
        print(f"[SPLIT] {page_info} 文本分段: {len(text)}字 → {len(segments)}段 (每段≤{MAX_SEGMENT_LENGTH}字)")

    # 如果只有一段，使用原来的逻辑
    if len(segments) == 1:
        return await _generate_single_segment(
            segments[0], output_file, voice_name, prompt_wav, prompt_text,
            registered_spk_id, available_urls, page_info
        )

    # 🆕 多段并发处理（先满血使用所有实例，故障时才转移）
    temp_dir = os.path.dirname(output_file)
    segment_files = []

    # 创建临时文件路径
    for i in range(len(segments)):
        seg_file = os.path.join(temp_dir, f"segment_{os.path.basename(output_file)}_{i}.wav")
        segment_files.append(seg_file)

    # 🔧 全局故障转移：区分"繁忙"与"真实故障"
    # 只记录真实故障（HTTP错误），繁忙（超时）不永久标记
    import asyncio
    failed_instances = set()  # 记录真实失败的实例
    busy_instances = set()    # 临时记录当前繁忙的实例
    failover_lock = asyncio.Lock()  # 保护失败集合的锁

    async def _generate_segment(idx):
        nonlocal failed_instances, busy_instances

        seg_text = segments[idx]
        seg_file = segment_files[idx]
        seg_page_info = f"{page_info}-段{idx+1}/{len(segments)}"

        # 轮询分配主实例（确保每个段使用不同的主实例）
        primary_url = available_urls[idx % len(available_urls)]

        # 先尝试主实例（增加超时到120秒）
        try:
            result = await _try_single_instance(
                seg_text, seg_file, voice_name, prompt_wav, prompt_text,
                registered_spk_id, primary_url, seg_page_info, timeout=120
            )
            if result:
                # 成功后从繁忙列表移除（如果之前被标记）
                async with failover_lock:
                    if primary_url in busy_instances:
                        busy_instances.discard(primary_url)
                return True
        except asyncio.TimeoutError:
            # 超时 = 繁忙，临时标记但不永久禁用
            async with failover_lock:
                if primary_url not in busy_instances:
                    busy_instances.add(primary_url)
                    print(f"[BUSY] {seg_page_info} 主实例 {primary_url} 繁忙（超时120秒），尝试备用实例...")
        except Exception as e:
            # HTTP错误等真实故障，永久标记
            error_msg = str(e)
            print(f"[FAIL] {seg_page_info} 主实例 {primary_url} 失败: {error_msg[:100]}")

            # 判断是否为真实故障（HTTP 4xx/5xx错误）
            is_real_failure = any(code in error_msg for code in ['400', '422', '500', '502', '503'])

            if is_real_failure:
                async with failover_lock:
                    if primary_url not in failed_instances:
                        failed_instances.add(primary_url)
                        print(f"[FAILOVER] {seg_page_info} 主实例 {primary_url} 标记为故障（HTTP错误）")

        # 准备备用实例列表（跳过真实故障，但考虑繁忙实例）
        available_backup = [url for url in available_urls if url not in failed_instances]

        if not available_backup:
            print(f"[ERROR] {seg_page_info} 所有实例均已故障")
            return False

        for backup_url in available_backup:
            # 跳过当前正在处理的主实例（避免重复）
            if backup_url == primary_url:
                continue

            try:
                print(f"[RETRY] {seg_page_info} 尝试备用实例 {backup_url}")
                result = await _try_single_instance(
                    seg_text, seg_file, voice_name, prompt_wav, prompt_text,
                    registered_spk_id, backup_url, seg_page_info, timeout=120
                )
                if result:
                    print(f"[SUCCESS] {seg_page_info} 备用实例 {backup_url} 成功!")
                    return True
            except asyncio.TimeoutError:
                # 备用实例也繁忙，不标记为故障
                print(f"[BUSY] {seg_page_info} 备用实例 {backup_url} 繁忙")
                continue
            except Exception as e:
                # 判断是否为真实故障
                error_msg = str(e)
                print(f"[FAIL] {seg_page_info} 备用实例 {backup_url} 失败: {error_msg[:100]}")

                is_real_failure = any(code in error_msg for code in ['400', '422', '500', '502', '503'])
                if is_real_failure:
                    async with failover_lock:
                        if backup_url not in failed_instances:
                            failed_instances.add(backup_url)
                continue

        return False

    # 🔧 工作池模式：始终保持 max_concurrent 个任务并发运行
    # 一个完成，立即启动下一个，充分利用所有实例
    max_concurrent = len(available_urls)  # 并发数 = 实例数
    semaphore = asyncio.Semaphore(max_concurrent)

    async def process_with_limit(idx):
        """使用信号量限制并发数"""
        async with semaphore:
            return await _generate_segment(idx)

    print(f"[POOL] 启动工作池：并发数={max_concurrent}，总分段数={len(segments)}")

    # 启动所有任务，但信号量会限制同时运行的只有 max_concurrent 个
    results = await asyncio.gather(*[process_with_limit(i) for i in range(len(segments))])

    # 检查是否所有分段都成功
    if not all(results):
        # 清理临时文件
        for f in segment_files:
            if os.path.exists(f): os.remove(f)
        return False

    # 🆕 使用 ffmpeg 合并音频段
    try:
        # 创建合并列表文件
        list_file = os.path.join(temp_dir, f"merge_list_{os.path.basename(output_file)}.txt")
        with open(list_file, "w", encoding="utf-8") as f:
            for seg_file in segment_files:
                f.write(f"file '{os.path.abspath(seg_file).replace(os.sep, '/')}'\n")

        # 使用 ffmpeg concat 合并音频
        # 注意：输入是 WAV 格式，输出也用 WAV 以避免重新编码
        # 如果输出文件名是 .mp3，ffmpeg 会自动转码
        cmd = [
            "ffmpeg", "-y", "-hide_banner", "-loglevel", "error",
            "-f", "concat", "-safe", "0", "-i", list_file,
            output_file
        ]
        process = await asyncio.create_subprocess_exec(
            *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await process.communicate()

        if process.returncode != 0:
            print(f"[ERROR] {page_info} 音频合并失败: {stderr.decode('utf-8')}")
            return False

        # 清理临时文件
        for f in segment_files:
            if os.path.exists(f): os.remove(f)
        if os.path.exists(list_file): os.remove(list_file)

        print(f"[OK] {page_info} 分段合并完成: {len(segments)}段 → {output_file}")
        return True

    except Exception as e:
        print(f"[ERROR] {page_info} 音频合并异常: {e}")
        return False


async def _try_single_instance(text, output_file, voice_name, prompt_wav, prompt_text,
                                registered_spk_id, base_url, page_info, timeout=60):
    """
    尝试单个实例生成音频（不进行故障转移）
    - base_url: 单个实例的URL
    - timeout: 超时时间（秒）
    - 返回: True成功, False失败
    - 抛出异常以便调用者处理故障转移
    """
    url = f"{base_url}/api/tts/sft"
    data = aiohttp.FormData()
    data.add_field("tts_text", text)
    data.add_field("speed", "1.0")

    file_obj = None
    try:
        if registered_spk_id:
            url = f"{base_url}/api/tts/zero_shot"
            data.add_field("speaker_id", registered_spk_id)
        elif voice_name == "zero_shot" and prompt_wav:
            url = f"{base_url}/api/tts/zero_shot"
            data.add_field("prompt_text", prompt_text)
            file_obj = open(prompt_wav, "rb")
            data.add_field("prompt_wav", file_obj,
                           filename=os.path.basename(prompt_wav),
                           content_type="audio/wav")
        else:
            data.add_field("speaker_id", voice_name)

        try:
            client_timeout = aiohttp.ClientTimeout(total=timeout)
            async with aiohttp.ClientSession(timeout=client_timeout) as session:
                async with session.post(url, data=data) as resp:
                    resp.raise_for_status()
                    content = await resp.read()
                    with open(output_file, "wb") as f:
                        f.write(content)
                    return True

        except asyncio.TimeoutError:
            raise Exception(f"超时({timeout}秒)")
        except aiohttp.ClientResponseError as e:
            raise Exception(f"HTTP错误({e.status})")
        except aiohttp.ClientConnectionError as e:
            raise Exception(f"连接失败")
        except Exception as e:
            raise Exception(f"未知错误: {str(e)[:50]}")

    finally:
        if file_obj and not file_obj.closed:
            file_obj.close()


async def _generate_single_segment(text, output_file, voice_name, prompt_wav, prompt_text,
                                    registered_spk_id, available_urls, page_info):
    """
    生成单个文本段的音频（内部函数，带自动故障转移）
    - available_urls: 可用实例URL列表，自动按顺序尝试
    - 用于单段文本场景
    """
    max_timeout = 60  # 单次尝试超时时间（秒）

    for url_idx, base_url in enumerate(available_urls):
        print(f"[TRY] {page_info} 尝试实例 {url_idx+1}/{len(available_urls)}: {base_url}")
        try:
            result = await _try_single_instance(
                text, output_file, voice_name, prompt_wav, prompt_text,
                registered_spk_id, base_url, page_info, timeout=max_timeout
            )
            if result:
                if url_idx > 0:
                    print(f"[FAILOVER_SUCCESS] {page_info} 故障转移成功! (实例{url_idx+1})")
                return True
        except Exception as e:
            print(f"[FAIL] {page_info} 实例 {base_url} 失败: {str(e)}")
            if url_idx < len(available_urls) - 1:
                print(f"  → 尝试下一个实例...")
                continue
            else:
                print(f"[ERROR] {page_info} 所有实例均已失败")
                return False

    return False

async def text_to_speech_wrapper(text, output_file, semaphore, voice_name,
                                  prompt_wav=None, prompt_text="",
                                  registered_spk_id=None, api_url=None, page_num=None, api_urls=None):
    async with semaphore:
        if not text.strip(): return True
        t0 = _time.time()

        # 🔧 自动检测音色类型：Neural 格式音色（Edge TTS）不走 CosyVoice
        is_neural_voice = isinstance(voice_name, str) and "Neural" in voice_name

        if is_neural_voice:
            # Edge TTS 格式（如 zh-CN-XiaoxiaoNeural）→ 走 Edge TTS
            result = await _generate_edge(text, output_file, voice_name)
        elif TTS_PROVIDER == "cosyvoice":
            result = await _generate_cosyvoice(text, output_file, voice_name,
                                               prompt_wav, prompt_text,
                                               registered_spk_id, api_url, page_num, api_urls)
        elif TTS_PROVIDER == "azure":
            result = await _generate_azure(text, output_file, voice_name)
        else:
            result = await _generate_edge(text, output_file, voice_name)
        elapsed = _time.time() - t0
        server_tag = api_url.split(':')[-1] if api_url else ''
        status = "[OK]" if result else "[ERROR]"
        # 4. 显示进度详情：显示页码
        page_info = f"第{page_num}页" if page_num else "N/A"
        print(f"  {status} [TTS] {page_info} | {os.path.basename(output_file)} | {elapsed:.1f}s | {len(text)}字 | :{server_tag}")
        return result

async def create_silent_audio(duration, output_path):
    if os.path.exists(output_path): return
    cmd = ["ffmpeg", "-y", "-hide_banner", "-loglevel", "error", "-f", "lavfi", "-i", "anullsrc=r=24000:cl=mono", "-t", str(duration), "-c:a", "libmp3lame", "-q:a", "4", output_path]
    subprocess.run(cmd, check=True)

def get_audio_duration(audio_path):
    try:
        cmd = ["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "json", audio_path]
        result = subprocess.run(cmd, capture_output=True, text=True, encoding="utf-8")
        return float(json.loads(result.stdout)['format']['duration'])
    except: return 3.0

def build_random_filter(duration):
    if duration < 2.0: return f"fade=t=in:st=0:d=0.5,fade=t=out:st={duration-0.5}:d=0.5", "Fade"
    effects = ["fade", "blur"]
    chosen = random.choice(effects)
    if chosen == "fade": vf = f"fade=t=in:st=0:d=0.5,fade=t=out:st={duration-0.5}:d=0.5"
    else: vf = f"boxblur=luma_radius=20:luma_power=1:enable='between(t,0,0.5)+between(t,{duration-0.5},{duration})',fade=t=in:st=0:d=0.5,fade=t=out:st={duration-0.5}:d=0.5"
    return vf, chosen

# --- 4. 渲染单页 (支持多模式) ---
async def render_slide_video(img_path, audio_path, output_video_path, video_mode="studio", effect_override=None):
    if os.path.exists(output_video_path): os.remove(output_video_path)
    duration = get_audio_duration(audio_path)
    raw_effect_filter, _ = build_random_filter(duration)

    cmd = []
    
    # 🌟 模式一：演播室模式 (叠加背景)
    if video_mode == "studio":
        if not os.path.exists(BACKGROUND_IMAGE_PATH): return None
        w, h = SCREEN_LAYOUT['w'], SCREEN_LAYOUT['h']
        x, y = SCREEN_LAYOUT['x'], SCREEN_LAYOUT['y']
        
        # 复杂滤镜链
        filter_complex = (
            f"[1:v]scale={w}:{h},setsar=1,{raw_effect_filter}[ppt];"
            f"[0:v][ppt]overlay=x={x}:y={y}:shortest=1[outv]"
        )
        cmd = [
            "ffmpeg", "-y", "-hide_banner", "-loglevel", "error",
            "-loop", "1", "-i", BACKGROUND_IMAGE_PATH,
            "-loop", "1", "-i", img_path,
            "-i", audio_path,
            "-filter_complex", filter_complex,
            "-map", "[outv]", "-map", "2:a",
            "-c:v", "h264_nvenc", "-preset", "p1", "-r", "24", "-pix_fmt", "yuv420p", "-shortest",
            output_video_path
        ]
        
    # 🌟 模式二：默认模式 (全屏PPT，无背景)
    else:
        cmd = [
            "ffmpeg", "-y", "-hide_banner", "-loglevel", "error",
            "-loop", "1", "-i", img_path, # 只有一个视频输入
            "-i", audio_path,
            "-vf", raw_effect_filter,     # 直接应用转场滤镜
            "-c:v", "h264_nvenc", "-preset", "p1", "-r", "24", "-pix_fmt", "yuv420p", "-shortest",
            output_video_path
        ]

    try:
        process = await asyncio.create_subprocess_exec(*cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE)
        # 使用 communicate() 读取输出，防止缓冲区填满导致的死锁 (假死)
        stdout, stderr = await process.communicate()

        if process.returncode != 0:
            print(f"[WARN] [FFmpeg Error] {stderr.decode('utf-8')}")
            return None

        return output_video_path
    except Exception as e:
        print(f"[ERROR] [Render Error] {e}")
        return None

# --- 主任务 (带进度回调) ---
async def generate_video_task(ppt_path, output_video_path, temp_dir, voice_name, video_mode, session_id):
    total_slides = 0

    # 动态发现可用实例（包含真实API测试）
    COSYVOICE_API_URLS = await _discover_cosyvoice_instances()

    # 根据实例数量动态调整并发数
    # 原则：每个实例分配1个并发请求，充分利用所有实例
    instance_count = len(COSYVOICE_API_URLS)
    max_tts_concurrent = max(instance_count, 1)  # 至少为1

    print(f"[TASK] ═════════════════════════════════════════")
    print(f"[TASK] CosyVoice 多实例并发模式")
    print(f"[TASK] 最终可用实例: {instance_count} 个")
    print(f"[TASK] 实例列表: {COSYVOICE_API_URLS}")
    print(f"[TASK] 客户端并发: {max_tts_concurrent} (每个实例1并发)")
    print(f"[TASK] ═════════════════════════════════════════")

    # ── 阶段 1：解析 PPT ──
    update_progress(session_id, "parse", 0, 0, "正在解析 PPT 提取幻灯片...")
    img_dir, vid_dir = os.path.join(temp_dir, "images"), os.path.join(temp_dir, "videos")
    if not os.path.exists(vid_dir): os.makedirs(vid_dir)
    if not ppt_to_images(ppt_path, img_dir):
        update_progress(session_id, "error", done=True, success=False, detail="PPT 解析失败")
        return False

    prs = Presentation(ppt_path)
    tts_tasks, slides_data = [], []
    tts_semaphore = asyncio.Semaphore(max_tts_concurrent)

    # 提取零样本克隆参数
    prompt_wav = voice_name.get("prompt_wav") if isinstance(voice_name, dict) else None
    prompt_text = voice_name.get("prompt_text", "") if isinstance(voice_name, dict) else ""
    real_voice_name = voice_name.get("voice_name", "中文女") if isinstance(voice_name, dict) else voice_name

    # 🆕 检查是否使用已保存的用户音色
    registered_spk_id = voice_name.get("registered_speaker_id") if isinstance(voice_name, dict) else None

    if not registered_spk_id and real_voice_name == "zero_shot" and prompt_wav:
        # 临时克隆模式：需要预注册到所有动态发现的实例
        update_progress(session_id, "parse", 0, 0, "正在预注册教师音色...")
        registered_spk_id = await _register_zero_shot_speaker(session_id, prompt_wav, prompt_text, COSYVOICE_API_URLS)

    # 🔧 重要：如果使用已保存的音色ID，需要确保它在所有实例上都存在
    # 因为新实例（如118服务器的实例）可能没有这个音色
    if registered_spk_id:
        # 验证音色是否在所有实例上可用
        is_available = await _verify_speaker_on_all_instances(registered_spk_id, COSYVOICE_API_URLS)

        # 如果不可用且有原始音频文件，重新分发到所有实例
        if not is_available and prompt_wav:
            print(f"[SYNC] 音色 {registered_spk_id} 不在所有实例上，正在同步...")
            await _register_zero_shot_speaker(session_id, prompt_wav, prompt_text, COSYVOICE_API_URLS)
            print(f"[SYNC] 音色同步完成")
        elif not is_available and not prompt_wav:
            print(f"[WARN] 音色 {registered_spk_id} 不在所有实例上，且缺少原始音频文件，部分实例可能失败")
            print(f"[WARN] 建议重新上传音频文件以确保一致性")

    print(f"[ENGINE] [Engine] 开始处理 | 模式: {video_mode} | 音色: {real_voice_name}" +
          (f" | speaker_id: {registered_spk_id}" if registered_spk_id else ""))

    for i, slide in enumerate(prs.slides):
        idx = i + 1
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide and slide.notes_slide.notes_text_frame else ""
        notes = notes.replace('\n', '，').strip()
        img, aud, vid = os.path.join(img_dir, f"{idx}.png"), os.path.join(temp_dir, f"audio_{idx}.wav"), os.path.join(vid_dir, f"seg_{idx}.mp4")
        if not os.path.exists(img): continue
        slides_data.append({"img": img, "aud": aud, "vid": vid, "notes": notes})
        if not notes:
            await create_silent_audio(3, aud)

    total_slides = len(slides_data)
    update_progress(session_id, "parse", total_slides, total_slides, f"解析完毕, 共 {total_slides} 页幻灯片")

    # ── 阶段 2：语音合成（按分段并发，而非按页） ──
    tts_done = 0
    tts_total_start = _time.time()
    num_servers = len(COSYVOICE_API_URLS)

    # 收集需要合成的页面
    tts_items = [(i, d) for i, d in enumerate(slides_data) if d["notes"]]
    tts_need = len(tts_items)

    if tts_items:
        update_progress(session_id, "tts", 0, total_slides,
                        f"正在合成语音，共 {tts_need} 页...")

        # 🆕 第一步：预分段，收集所有分段任务
        MAX_SEGMENT_LENGTH = 200  # 每段最大字符数
        all_segments = []  # (page_idx, page_num, segment_idx, segment_text)
        for idx, d in tts_items:
            text = d["notes"]
            segments = _split_long_text(text, MAX_SEGMENT_LENGTH)
            page_num = idx + 1
            for seg_idx, seg_text in enumerate(segments):
                all_segments.append((idx, page_num, seg_idx, seg_text, d))

        total_segments = len(all_segments)
        print(f"[TTS] 预分段完成: {tts_need}页 → {total_segments}段")

        # 🔧 第二步：按分段并发（始终保持 max_concurrent 个分段在处理）
        max_concurrent = num_servers  # 并发数 = 实例数
        segment_semaphore = asyncio.Semaphore(max_concurrent)

        tts_failed = []  # 记录失败的页
        completed_pages = set()  # 记录已完成的页

        async def _do_segment_tts(page_idx, page_num, seg_idx, seg_text, slide_data):
            """处理单个分段（带故障转移）"""
            nonlocal tts_done

            async with segment_semaphore:
                # 生成分段音频临时文件
                temp_dir = os.path.dirname(slide_data["aud"])
                seg_file = os.path.join(temp_dir, f"seg_p{page_num}_s{seg_idx}.wav")

                seg_page_info = f"第{page_num}页-段{seg_idx+1}"

                # 🌟 改进：使用全局轮询计数器，确保任务均匀分布到所有实例
                # 问题：之前的哈希算法 (page_idx * 10 + seg_idx) % 5 导致所有任务都落到第一个实例
                # 解决：使用全局计数器实现真正的轮询（Round-Robin）
                global _tts_round_robin_counter
                try:
                    _tts_round_robin_counter
                except NameError:
                    _tts_round_robin_counter = 0

                # 使用全局计数器获取下一个实例索引（确保均匀分布）
                list_len = len(COSYVOICE_API_URLS)
                primary_idx = _tts_round_robin_counter % list_len
                _tts_round_robin_counter += 1

                # 构建轮询队列：[当前选中的实例, 其他实例按顺序]
                rotated_urls = [COSYVOICE_API_URLS[primary_idx]] + \
                               [COSYVOICE_API_URLS[i] for i in range(list_len) if i != primary_idx]

                # 🔧 Neural 音色（Edge TTS）直接走 Edge TTS，不走 CosyVoice
                is_neural = isinstance(real_voice_name, str) and "Neural" in real_voice_name
                if is_neural:
                    print(f"[SEND] {seg_page_info} -> Edge TTS ({real_voice_name})")
                    result = await _generate_edge(seg_text, seg_file, real_voice_name)
                    if result:
                        return page_idx, page_num, seg_idx, seg_file, True
                    print(f"[FAIL] {seg_page_info} Edge TTS 失败")
                    return page_idx, page_num, seg_idx, seg_file, False

                # 🔧 故障转移：遍历轮询好的实例队列（CosyVoice 路径）
                for attempt_idx, api_url in enumerate(rotated_urls):
                    try:
                        if attempt_idx == 0:
                            print(f"[SEND] {seg_page_info} -> {api_url}")
                        
                        result = await _try_single_instance(
                            seg_text, seg_file, real_voice_name, prompt_wav, prompt_text,
                            registered_spk_id, api_url, seg_page_info, timeout=300
                        )
                        # 成功
                        if attempt_idx > 0:
                            print(f"[RETRY_SUCCESS] {seg_page_info} 备用实例 {api_url} 成功")
                        else:
                            print(f"[SUCCESS] {seg_page_info} 实例 {api_url} 生成完毕!")
                        return page_idx, page_num, seg_idx, seg_file, True
                    except Exception as e:
                        # 失败，尝试下一个实例
                        print(f"[FAIL] {seg_page_info} 实例 {api_url} 失败: {str(e)[:80]}")
                        continue

                # 所有实例都失败，输出文本内容用于调试
                print(f"[ERROR] {seg_page_info} 所有实例均失败")
                if len(seg_text) > 100:
                    print(f"[DEBUG] 失败文本({len(seg_text)}字): {seg_text[:100]}...")
                else:
                    print(f"[DEBUG] 失败文本({len(seg_text)}字): {seg_text}")
                return page_idx, page_num, seg_idx, seg_file, False

        # 启动所有分段任务（信号量会限制并发数）
        segment_coros = [
            _do_segment_tts(idx, page_num, seg_idx, seg_text, d)
            for idx, page_num, seg_idx, seg_text, d in all_segments
        ]
        segment_results = await asyncio.gather(*segment_coros)

        # 🆕 第三步：合并每个页面的分段音频
        from collections import defaultdict
        page_segments = defaultdict(list)  # page_idx -> [(seg_idx, seg_file), ...]

        for page_idx, page_num, seg_idx, seg_file, success in segment_results:
            if success:
                page_segments[page_idx].append((seg_idx, seg_file))
            else:
                # 标记页面失败
                if page_num not in tts_failed:
                    tts_failed.append(page_num)

        # 合并每个页面的分段
        for page_idx, segments_list in page_segments.items():
            if not segments_list:
                continue

            # 按seg_idx排序
            segments_list.sort(key=lambda x: x[0])

            # 如果只有一个分段，直接重命名
            if len(segments_list) == 1:
                _, seg_file = segments_list[0]
                target_file = slides_data[page_idx]["aud"]
                if os.path.exists(seg_file):
                    os.rename(seg_file, target_file)
                    tts_done += 1
                    update_progress(session_id, "tts", tts_done, total_slides,
                                    f"已完成 {tts_done}/{tts_need} 页语音合成")
            else:
                # 多个分段，使用ffmpeg合并
                target_file = slides_data[page_idx]["aud"]
                temp_dir = os.path.dirname(target_file)
                list_file = os.path.join(temp_dir, f"merge_p{page_idx}.txt")

                with open(list_file, "w", encoding="utf-8") as f:
                    for seg_idx, seg_file in segments_list:
                        f.write(f"file '{os.path.abspath(seg_file).replace(os.sep, '/')}'\n")

                cmd = [
                    "ffmpeg", "-y", "-hide_banner", "-loglevel", "error",
                    "-f", "concat", "-safe", "0", "-i", list_file,
                    target_file
                ]

                try:
                    process = await asyncio.create_subprocess_exec(*cmd,
                        stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE)
                    await process.communicate()

                    if process.returncode == 0:
                        # 清理临时分段文件
                        for _, seg_file in segments_list:
                            if os.path.exists(seg_file):
                                os.remove(seg_file)
                        if os.path.exists(list_file):
                            os.remove(list_file)

                        tts_done += 1
                        update_progress(session_id, "tts", tts_done, total_slides,
                                        f"已完成 {tts_done}/{tts_need} 页语音合成")
                    else:
                        if page_idx + 1 not in tts_failed:
                            tts_failed.append(page_idx + 1)
                except Exception as e:
                    print(f"[ERROR] 合并第{page_idx+1}页音频失败: {e}")
                    if page_idx + 1 not in tts_failed:
                        tts_failed.append(page_idx + 1)

        if tts_failed:
            update_progress(session_id, "error", done=True, success=False,
                            detail=f"第 {tts_failed[0]} 页语音合成失败")
            if registered_spk_id: await _unregister_speaker(registered_spk_id)
            return False

        update_progress(session_id, "tts", total_slides, total_slides,
                        f"全部 {tts_need} 页语音合成完成")
    else:
        tts_done = total_slides


    tts_total_elapsed = _time.time() - tts_total_start
    print(f"[TIME] [TTS 总耗时] {tts_total_elapsed:.1f}s | {len(tts_items)} 页 | {num_servers} 实例并行 | 平均 {tts_total_elapsed/max(len(tts_items),1):.1f}s/页")

    # ── 阶段 3：视频渲染 ──
    render_done = 0
    render_sem = asyncio.Semaphore(MAX_RENDER_CONCURRENT)

    async def do_render(idx, d):
        nonlocal render_done
        async with render_sem:
            if not os.path.exists(d['aud']): return None
            result = await render_slide_video(d['img'], d['aud'], d['vid'], video_mode=video_mode)
            render_done += 1
            update_progress(session_id, "render", render_done, total_slides,
                            f"已渲染 {render_done}/{total_slides} 页视频")
            return result

    update_progress(session_id, "render", 0, total_slides, "正在渲染视频片段...")
    render_tasks = [do_render(i, d) for i, d in enumerate(slides_data)]
    valid_vids = [v for v in await asyncio.gather(*render_tasks) if v]
    if not valid_vids:
        update_progress(session_id, "error", done=True, success=False, detail="视频渲染失败")
        return False

    # ── 阶段 4：合并输出 ──
    update_progress(session_id, "merge", 0, 1, "正在合并视频片段为最终文件...")
    list_path = os.path.join(temp_dir, "list.txt")
    with open(list_path, "w", encoding="utf-8") as f:
        for v in valid_vids: f.write(f"file '{os.path.abspath(v).replace(os.sep, '/')}'\n")

    subprocess.run(["ffmpeg", "-y", "-hide_banner", "-loglevel", "error", "-f", "concat", "-safe", "0", "-i", list_path, "-c", "copy", output_video_path])
    cleanup_folder(temp_dir)
    # 只清理临时音色（p2v_ 前缀），不清理用户已保存的音色（u 前缀）
    if registered_spk_id and registered_spk_id.startswith("p2v_"):
        await _unregister_speaker(registered_spk_id)
    print(f"[OK] 完成: {output_video_path}")

    update_progress(session_id, "done", 1, 1, "视频生成完成！", done=True, success=True)
    return True

# 🆕 添加入口参数 video_mode
def run_generation(ppt_path, output_path, session_id, voice_name, video_mode="studio", effect_type="random"):
    temp_dir = os.path.join(os.path.dirname(output_path), f"temp_{session_id}")
    update_progress(session_id, "init", 0, 0, "任务已提交，正在初始化...")
    try:
        asyncio.run(generate_video_task(ppt_path, output_path, temp_dir, voice_name, video_mode, session_id))
        return True
    except Exception as e:
        print(f"[ERROR] 错误: {e}")
        update_progress(session_id, "error", done=True, success=False, detail=f"系统错误: {e}")
        cleanup_folder(temp_dir)
        return False